"""Bounded, native-only document validation and extraction."""

from __future__ import annotations

import csv
import hashlib
import io
import re
import unicodedata
import zipfile
from dataclasses import dataclass
from pathlib import Path

import bleach
import filetype
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from PIL import Image, UnidentifiedImageError
from pptx import Presentation
from pypdf import PdfReader

from akc_api.settings import Settings

ALLOWED_EXTENSIONS = {
    ".pdf",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".tif",
    ".tiff",
    ".docx",
    ".pptx",
    ".xlsx",
    ".csv",
    ".html",
    ".htm",
    ".txt",
    ".md",
    ".vtt",
    ".srt",
}
TEXT_EXTENSIONS = {".csv", ".html", ".htm", ".txt", ".md", ".vtt", ".srt"}
ZIP_EXTENSIONS = {".docx", ".pptx", ".xlsx"}
MIME_BY_EXTENSION = {
    ".pdf": {"application/pdf"},
    ".png": {"image/png"},
    ".jpg": {"image/jpeg"},
    ".jpeg": {"image/jpeg"},
    ".webp": {"image/webp"},
    ".tif": {"image/tiff"},
    ".tiff": {"image/tiff"},
    ".docx": {"application/vnd.openxmlformats-officedocument.wordprocessingml.document"},
    ".pptx": {"application/vnd.openxmlformats-officedocument.presentationml.presentation"},
    ".xlsx": {"application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"},
    ".csv": {"text/csv", "text/plain", "application/csv"},
    ".html": {"text/html", "application/xhtml+xml"},
    ".htm": {"text/html", "application/xhtml+xml"},
    ".txt": {"text/plain"},
    ".md": {"text/markdown", "text/plain"},
    ".vtt": {"text/vtt", "text/plain"},
    ".srt": {"application/x-subrip", "text/plain"},
}


class FileValidationError(ValueError):
    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code


@dataclass(frozen=True)
class ParsedPage:
    page_number: int
    text: str
    width_pt: float | None = None
    height_pt: float | None = None
    image_coverage: float = 0.0


@dataclass(frozen=True)
class ParsedDocument:
    document_type: str
    pages: tuple[ParsedPage, ...]


def safe_filename(filename: str) -> str:
    name = unicodedata.normalize("NFC", Path(filename).name)
    name = re.sub(r'[<>:"/\\|?*\x00-\x1f]', "-", name)
    name = re.sub(r"\s+", " ", name).strip(" .")
    if not name:
        name = "upload"
    return name[:240]


def _check_archive(data: bytes, settings: Settings) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            infos = archive.infolist()
            if len(infos) > settings.max_archive_files:
                raise FileValidationError("ARCHIVE_FILE_LIMIT", "archive has too many files")
            total = sum(info.file_size for info in infos)
            compressed = max(1, sum(info.compress_size for info in infos))
            if total > settings.max_archive_uncompressed_bytes:
                raise FileValidationError(
                    "ARCHIVE_SIZE_LIMIT", "archive expands beyond the configured limit"
                )
            if total / compressed > settings.max_archive_ratio:
                raise FileValidationError(
                    "ARCHIVE_RATIO_LIMIT", "archive compression ratio is unsafe"
                )
            for info in infos:
                parts = Path(info.filename).parts
                if info.filename.startswith(("/", "\\")) or ".." in parts:
                    raise FileValidationError("ARCHIVE_PATH_TRAVERSAL", "unsafe archive path")
            relationship_files = [item for item in infos if item.filename.endswith(".rels")]
            for info in relationship_files:
                content = archive.read(info)
                if b'TargetMode="External"' in content:
                    raise FileValidationError(
                        "OFFICE_EXTERNAL_RELATION",
                        "external Office relationships are not allowed",
                    )
    except zipfile.BadZipFile as exc:
        raise FileValidationError("INVALID_OFFICE_ARCHIVE", "invalid Office archive") from exc


def validate_file(
    *, filename: str, declared_mime: str, data: bytes, expected_sha256: str, settings: Settings
) -> tuple[str, str]:
    extension = Path(filename).suffix.lower()
    if extension not in ALLOWED_EXTENSIONS:
        raise FileValidationError("UNSUPPORTED_FILE", "file extension is not allowed")
    if len(data) > settings.max_upload_bytes:
        raise FileValidationError("FILE_TOO_LARGE", "file exceeds the configured size limit")
    digest = hashlib.sha256(data).hexdigest()
    if digest.lower() != expected_sha256.lower():
        raise FileValidationError("CHECKSUM_MISMATCH", "uploaded bytes do not match SHA-256")
    declared_base = declared_mime.split(";", 1)[0].strip().lower()
    allowed = MIME_BY_EXTENSION[extension]
    if declared_base not in allowed:
        raise FileValidationError("MIME_MISMATCH", "declared MIME does not match extension")
    kind = filetype.guess(data[:8192])
    if extension == ".pdf" and (kind is None or kind.mime != "application/pdf"):
        raise FileValidationError("MAGIC_MISMATCH", "PDF magic bytes are invalid")
    if extension in ZIP_EXTENSIONS:
        if not data.startswith(b"PK"):
            raise FileValidationError("MAGIC_MISMATCH", "Office magic bytes are invalid")
        _check_archive(data, settings)
    if extension in TEXT_EXTENSIONS and b"\x00" in data[:8192]:
        raise FileValidationError("BINARY_TEXT_FILE", "text source contains NUL bytes")
    return extension, digest


def _decode_text(data: bytes) -> str:
    for encoding in ("utf-8-sig", "utf-8", "cp949", "utf-16"):
        try:
            return data.decode(encoding)
        except UnicodeDecodeError:
            continue
    raise FileValidationError("TEXT_ENCODING", "text encoding is unsupported")


def _parse_pdf(
    data: bytes,
    max_pages: int,
    *,
    password: bytes | None = None,
) -> ParsedDocument:
    reader = PdfReader(io.BytesIO(data), strict=True)
    if reader.is_encrypted:
        if password is None:
            raise FileValidationError("ENCRYPTED_PDF", "a PDF password is required")
        try:
            decrypted = bool(reader.decrypt(password))
        except Exception as exc:
            raise FileValidationError(
                "PDF_PASSWORD_INVALID", "the PDF password is invalid"
            ) from exc
        if not decrypted:
            raise FileValidationError("PDF_PASSWORD_INVALID", "the PDF password is invalid")
    if len(reader.pages) > max_pages:
        raise FileValidationError("PAGE_LIMIT", "document exceeds page limit")
    pages = []
    for number, page in enumerate(reader.pages, 1):
        box = page.mediabox
        pages.append(
            ParsedPage(
                page_number=number,
                text=page.extract_text() or "",
                width_pt=float(box.width),
                height_pt=float(box.height),
            )
        )
    return ParsedDocument("pdf", tuple(pages))


def _parse_docx(data: bytes) -> ParsedDocument:
    document = DocxDocument(io.BytesIO(data))
    lines = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            lines.append(" | ".join(cell.text for cell in row.cells))
    return ParsedDocument("docx", (ParsedPage(1, "\n".join(lines)),))


def _parse_pptx(data: bytes, max_pages: int) -> ParsedDocument:
    deck = Presentation(io.BytesIO(data))
    if len(deck.slides) > max_pages:
        raise FileValidationError("PAGE_LIMIT", "presentation exceeds page limit")
    pages = []
    for number, slide in enumerate(deck.slides, 1):
        lines = [
            shape.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text
        ]
        pages.append(ParsedPage(number, "\n".join(lines)))
    return ParsedDocument("pptx", tuple(pages))


def _parse_xlsx(data: bytes, max_pages: int) -> ParsedDocument:
    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    if len(workbook.worksheets) > max_pages:
        raise FileValidationError("PAGE_LIMIT", "workbook exceeds worksheet limit")
    pages = []
    for number, sheet in enumerate(workbook.worksheets, 1):
        lines = []
        for row in sheet.iter_rows(values_only=True):
            values = ["" if value is None else str(value) for value in row]
            if any(values):
                lines.append("\t".join(values))
        pages.append(ParsedPage(number, "\n".join(lines)))
    workbook.close()
    return ParsedDocument("xlsx", tuple(pages))


def _parse_csv(data: bytes) -> ParsedDocument:
    source = _decode_text(data)
    rows = csv.reader(io.StringIO(source))
    return ParsedDocument(
        "csv",
        (ParsedPage(1, "\n".join("\t".join(cell for cell in row) for row in rows)),),
    )


def _parse_html(data: bytes) -> ParsedDocument:
    source = _decode_text(data)
    clean = bleach.clean(
        source,
        tags=["p", "h1", "h2", "h3", "h4", "h5", "h6", "table", "tr", "th", "td"],
        attributes={},
        protocols=[],
        strip=True,
    )
    soup = BeautifulSoup(clean, "lxml")
    return ParsedDocument("html", (ParsedPage(1, soup.get_text("\n", strip=True)),))


def _parse_image(data: bytes, extension: str) -> ParsedDocument:
    try:
        with Image.open(io.BytesIO(data)) as image:
            image.verify()
        with Image.open(io.BytesIO(data)) as image:
            width, height = image.size
            if width <= 0 or height <= 0 or width * height > 100_000_000:
                raise FileValidationError("IMAGE_DIMENSIONS", "image dimensions are unsafe")
    except (Image.DecompressionBombError, UnidentifiedImageError, OSError) as exc:
        raise FileValidationError("INVALID_IMAGE", "image cannot be safely decoded") from exc
    return ParsedDocument(
        extension.removeprefix("."),
        (
            ParsedPage(
                page_number=1,
                text="",
                width_pt=float(width),
                height_pt=float(height),
                image_coverage=1.0,
            ),
        ),
    )


def parse_document(
    filename: str,
    data: bytes,
    settings: Settings,
    *,
    pdf_password: bytes | None = None,
) -> ParsedDocument:
    extension = Path(filename).suffix.lower()
    if extension == ".pdf":
        parsed = _parse_pdf(data, settings.max_pages, password=pdf_password)
    elif extension == ".docx":
        parsed = _parse_docx(data)
    elif extension == ".pptx":
        parsed = _parse_pptx(data, settings.max_pages)
    elif extension == ".xlsx":
        parsed = _parse_xlsx(data, settings.max_pages)
    elif extension == ".csv":
        parsed = _parse_csv(data)
    elif extension in {".html", ".htm"}:
        parsed = _parse_html(data)
    elif extension in {".png", ".jpg", ".jpeg", ".webp", ".tif", ".tiff"}:
        parsed = _parse_image(data, extension)
    else:
        parsed = ParsedDocument(extension.removeprefix("."), (ParsedPage(1, _decode_text(data)),))
    if len(parsed.pages) > settings.max_pages:
        raise FileValidationError("PAGE_LIMIT", "document exceeds page limit")
    return parsed


def page_preflight(text: str, *, image_coverage: float = 0.0) -> dict[str, float | int | bool]:
    replacement = text.count("\ufffd") / max(1, len(text))
    controls = sum(ord(char) < 32 and char not in "\n\r\t" for char in text)
    return {
        "native_text_chars": len(text),
        "native_words": len(text.split()),
        "invalid_char_ratio": controls / max(1, len(text)),
        "replacement_char_ratio": replacement,
        "image_coverage": image_coverage,
        "suspicious_text_layer": replacement > 0.001,
    }
