"""Upload-to-database evidence for structured native non-PDF parsing."""

from __future__ import annotations

import hashlib
import io
import uuid
from collections.abc import AsyncIterator
from dataclasses import dataclass
from pathlib import Path
from typing import Any

import httpx
import pytest_asyncio
from akc_api.artifacts import build_canonical_document
from akc_api.main import create_app
from akc_api.models import AnalysisTask, Block, Document, Export, Page
from akc_api.settings import Settings
from akc_worker_document.worker import AnalysisRuntime, AnalysisWorker
from docx import Document as WordDocument
from openpyxl import Workbook
from openpyxl.worksheet.table import Table as WorksheetTable
from pptx import Presentation
from pptx.util import Inches
from sqlalchemy import select

_TEST_SUPPORT_KEY = "native-nonpdf-worker-verification-key"
_REPOSITORY = Path(__file__).parents[3]
_FIXTURES = _REPOSITORY / "tests" / "fixtures" / "nonpdf"


@dataclass(frozen=True, slots=True)
class NativeSample:
    filename: str
    content_type: str
    payload: bytes
    document_type: str
    page_count: int
    native_prefix: str
    required_block_types: frozenset[str]


@pytest_asyncio.fixture
async def native_worker_api(
    tmp_path: Path,
) -> AsyncIterator[tuple[httpx.AsyncClient, Any, Settings]]:
    settings = Settings(
        env="test",
        database_url=f"sqlite+aiosqlite:///{(tmp_path / 'native-worker.db').as_posix()}",
        data_dir=tmp_path / "data",
        local_background_tasks=False,
        local_analysis_worker_enabled=False,
        analysis_max_source_bytes=2 * 1024 * 1024,
        analysis_max_attempts=2,
        analysis_lease_seconds=70,
        analysis_attempt_timeout_seconds=60,
        analysis_backoff_base_seconds=0.01,
        analysis_backoff_max_seconds=0.02,
        free_daily_file_cap=20,
        free_daily_page_cap=100,
        clamav_enabled=False,
        allow_development_antivirus_bypass=True,
        test_support_key=_TEST_SUPPORT_KEY,
    )
    app = create_app(settings)
    async with (
        app.router.lifespan_context(app),
        httpx.AsyncClient(
            transport=httpx.ASGITransport(app=app, raise_app_exceptions=True),
            base_url="http://testserver",
        ) as client,
    ):
        yield client, app, settings


async def _register(client: httpx.AsyncClient) -> dict[str, Any]:
    email = "native-worker@example.test"
    response = await client.post(
        "/v1/auth/register",
        json={
            "email": email,
            "password": "correct horse battery staple",
            "display_name": "Native Worker",
            "tenant_name": "Native Parser Evidence",
        },
    )
    assert response.status_code == 201, response.text
    captured = await client.post(
        "/__test__/verification-token",
        headers={"X-AKC-Test-Support-Key": _TEST_SUPPORT_KEY},
        json={"email": email},
    )
    assert captured.status_code == 200
    verified = await client.post(
        "/v1/auth/verify-email",
        json={"token": captured.json()["token"]},
    )
    assert verified.status_code == 200, verified.text
    return verified.json()


async def _upload(
    client: httpx.AsyncClient,
    *,
    project_id: str,
    sample: NativeSample,
) -> str:
    digest = hashlib.sha256(sample.payload).hexdigest()
    initiated = await client.post(
        "/v1/uploads/initiate",
        json={
            "project_id": project_id,
            "filename": sample.filename,
            "size": len(sample.payload),
            "content_type": sample.content_type,
            "sha256": digest,
        },
    )
    assert initiated.status_code == 201, initiated.text
    target = initiated.json()
    uploaded = await client.put(
        target["upload_url"],
        content=sample.payload,
        headers=target["headers"],
    )
    assert uploaded.status_code == 204, uploaded.text
    completed = await client.post(
        f"/v1/uploads/{target['upload_id']}/complete",
        json={"sha256": digest},
    )
    assert completed.status_code == 200, completed.text
    return str(target["document_id"])


async def _analyze(
    client: httpx.AsyncClient,
    worker: AnalysisWorker,
    document_id: str,
) -> uuid.UUID:
    queued = await client.post(f"/v1/documents/{document_id}/analyze")
    assert queued.status_code == 202, queued.text
    task_id = uuid.UUID(queued.json()["task_id"])
    assert await worker.run_once(task_id=task_id) is True
    return task_id


def _docx_bytes() -> bytes:
    document = WordDocument()
    document.core_properties.title = "구조 보존 DOCX"
    document.add_paragraph("구조 보존 DOCX", style="Title")
    document.add_heading("실험 결과", level=1)
    document.add_paragraph("원문 문단입니다.")
    document.add_paragraph("첫 번째 항목", style="List Bullet")
    table = document.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "구성"
    table.cell(0, 1).text = "점수"
    table.cell(0, 2).text = "비고"
    table.cell(1, 0).text = "개선"
    table.cell(1, 1).text = "0.94"
    table.cell(1, 1).merge(table.cell(1, 2))
    document.sections[0].header.paragraphs[0].text = "보안 등급: 내부"
    document.sections[0].footer.paragraphs[0].text = "문서 끝"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "PPTX 구조"
    body = first.placeholders[1].text_frame
    body.text = "첫 번째 문단"
    bullet = body.add_paragraph()
    bullet.text = "핵심 항목"
    bullet.level = 1
    table = first.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(4),
        Inches(4),
        Inches(1.5),
    ).table
    table.cell(0, 0).text = "구성"
    table.cell(0, 1).text = "점수"
    table.cell(1, 0).text = "기준"
    table.cell(1, 1).text = "0.86"
    first.notes_slide.notes_text_frame.text = "발표자 전용 메모"
    second = deck.slides.add_slide(deck.slide_layouts[5])
    second.shapes.title.text = "두 번째 슬라이드"
    second.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(5),
        Inches(1),
    ).text_frame.text = "슬라이드 순서를 보존합니다."
    output = io.BytesIO()
    deck.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Results"
    sheet.append(["구성", "점수", "비고", "병합", "합계"])
    sheet.append(["기준", 0.86, "검증 전", None, "=SUM(B2:B3)"])
    sheet.append(["개선", 0.94, "검증 후", None, None])
    sheet.merge_cells("C2:D2")
    sheet.row_dimensions[3].hidden = True
    sheet.add_table(WorksheetTable(displayName="EvidenceTable", ref="A1:B3"))
    hidden = workbook.create_sheet("Raw")
    hidden.sheet_state = "hidden"
    hidden.append(["raw", "value"])
    hidden.append(["alpha", 1])
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _samples() -> tuple[NativeSample, ...]:
    return (
        NativeSample(
            filename="evidence.docx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            ),
            payload=_docx_bytes(),
            document_type="docx",
            page_count=1,
            native_prefix="docx/",
            required_block_types=frozenset(
                {"title", "heading", "paragraph", "list", "table", "header", "footer"}
            ),
        ),
        NativeSample(
            filename="deck.pptx",
            content_type=(
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            ),
            payload=_pptx_bytes(),
            document_type="pptx",
            page_count=2,
            native_prefix="pptx/",
            required_block_types=frozenset(
                {"title", "heading", "paragraph", "list", "table", "footnote"}
            ),
        ),
        NativeSample(
            filename="workbook.xlsx",
            content_type=("application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"),
            payload=_xlsx_bytes(),
            document_type="xlsx",
            page_count=2,
            native_prefix="xlsx/",
            required_block_types=frozenset({"heading", "table"}),
        ),
        NativeSample(
            filename="article.html",
            content_type="text/html",
            payload=(_FIXTURES / "sample.html").read_bytes(),
            document_type="html",
            page_count=1,
            native_prefix="html/",
            required_block_types=frozenset(
                {"title", "heading", "paragraph", "list", "table", "figure", "caption"}
            ),
        ),
        NativeSample(
            filename="captions.srt",
            content_type="application/x-subrip",
            payload=(_FIXTURES / "sample.srt").read_bytes(),
            document_type="srt",
            page_count=1,
            native_prefix="srt/",
            required_block_types=frozenset({"title", "paragraph"}),
        ),
        NativeSample(
            filename="captions.vtt",
            content_type="text/vtt",
            payload=(_FIXTURES / "sample.vtt").read_bytes(),
            document_type="vtt",
            page_count=1,
            native_prefix="vtt/",
            required_block_types=frozenset({"title", "paragraph"}),
        ),
    )


async def test_six_native_formats_persist_structured_cir_and_provenance(
    native_worker_api: tuple[httpx.AsyncClient, Any, Settings],
) -> None:
    client, app, settings = native_worker_api
    registration = await _register(client)
    project = await client.post(
        "/v1/projects",
        json={"name": "Structured Native Matrix"},
    )
    assert project.status_code == 201, project.text
    project_id = str(project.json()["id"])
    worker = AnalysisWorker(
        engine=app.state.database.engine,
        store=app.state.object_store,
        runtime=AnalysisRuntime.from_api_settings(settings),
    )

    for sample in _samples():
        document_id = await _upload(
            client,
            project_id=project_id,
            sample=sample,
        )
        task_id = await _analyze(client, worker, document_id)
        async with app.state.database.sessions() as session:
            task = await session.get(AnalysisTask, task_id)
            document = await session.get(Document, uuid.UUID(document_id))
            pages = list(
                (
                    await session.scalars(
                        select(Page)
                        .where(Page.document_id == uuid.UUID(document_id))
                        .order_by(Page.page_number)
                    )
                ).all()
            )
            blocks = list(
                (
                    await session.scalars(
                        select(Block)
                        .where(Block.document_id == uuid.UUID(document_id))
                        .order_by(Block.block_order)
                    )
                ).all()
            )
            assert task is not None and task.status == "completed"
            assert task.last_error_code is None
            assert document is not None
            assert document.document_type == sample.document_type
            assert document.page_count == sample.page_count
            assert document.cir_schema_version == "cir-1.0.0"
            assert len(pages) == sample.page_count
            assert len(blocks) == task.block_count
            assert [block.block_order for block in blocks] == list(range(len(blocks)))
            assert sample.required_block_types.issubset({block.block_type for block in blocks})
            page_number_by_id = {page.id: page.page_number for page in pages}
            for block in blocks:
                assert block.page_id is not None
                structured = block.structured_content
                assert structured is not None
                assert structured["schemaVersion"] == "akc-native-block-1.0"
                source_refs = structured["sourceRefs"]
                assert source_refs
                assert source_refs[0]["nativeObjectId"].startswith(sample.native_prefix)
                assert source_refs[0]["pageNumber1"] == page_number_by_id[block.page_id]
                if source_refs[0].get("bbox1000") is not None:
                    assert block.bbox1000 == source_refs[0]["bbox1000"]
            first_native_structure = pages[0].preflight_metrics["native_structure"]
            assert first_native_structure["schemaVersion"] == "cir-1.0.0"
            assert (
                first_native_structure["documentMetadata"]["documentType"] == sample.document_type
            )

            export = Export(
                tenant_id=uuid.UUID(registration["tenant_id"]),
                project_id=uuid.UUID(project_id),
                document_id=uuid.UUID(document_id),
                export_type="portable",
                status="queued",
                options={},
                created_by=uuid.UUID(registration["user_id"]),
            )
            session.add(export)
            await session.flush()
            rebuilt, _knowledge = await build_canonical_document(session, export)
            assert len(rebuilt.blocks) == len(blocks)
            assert all(
                (source_ref.native_object_id or "").startswith(sample.native_prefix)
                for block in rebuilt.blocks
                for source_ref in block.source_refs
            )
            assert all(
                block.table is not None for block in rebuilt.blocks if block.type.value == "table"
            )


async def test_native_parser_error_code_reaches_analysis_task(
    native_worker_api: tuple[httpx.AsyncClient, Any, Settings],
) -> None:
    client, app, settings = native_worker_api
    await _register(client)
    project = await client.post(
        "/v1/projects",
        json={"name": "Native Failure"},
    )
    assert project.status_code == 201
    sample = NativeSample(
        filename="empty.vtt",
        content_type="text/vtt",
        payload=b"WEBVTT\n\nNOTE metadata only\n",
        document_type="vtt",
        page_count=1,
        native_prefix="vtt/",
        required_block_types=frozenset(),
    )
    document_id = await _upload(
        client,
        project_id=str(project.json()["id"]),
        sample=sample,
    )
    worker = AnalysisWorker(
        engine=app.state.database.engine,
        store=app.state.object_store,
        runtime=AnalysisRuntime.from_api_settings(settings),
    )
    task_id = await _analyze(client, worker, document_id)
    async with app.state.database.sessions() as session:
        task = await session.get(AnalysisTask, task_id)
        pages = list(
            (
                await session.scalars(
                    select(Page).where(Page.document_id == uuid.UUID(document_id))
                )
            ).all()
        )
        blocks = list(
            (
                await session.scalars(
                    select(Block).where(Block.document_id == uuid.UUID(document_id))
                )
            ).all()
        )
    assert task is not None
    assert task.status == "dead_letter"
    assert task.last_error_code == "VTT_NO_CUES"
    assert pages == []
    assert blocks == []
