"""PPTX slide, grouped reading-order, asset, chart, table, and note extraction."""

from __future__ import annotations

import io
import math
import zipfile
from dataclasses import dataclass
from typing import Any

from akc_cir import BlockType
from pptx import Presentation as open_presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER

from .models import (
    CirBuilder,
    SourceLocation,
    StructuredParseError,
    TableCellSpec,
    normalize_text,
)

_TITLE_PLACEHOLDERS = {
    PP_PLACEHOLDER.TITLE,
    PP_PLACEHOLDER.CENTER_TITLE,
}


@dataclass(frozen=True, slots=True)
class _ShapeNode:
    z_path: str
    shape: Any
    group_path: str | None
    container_top: int
    container_left: int


def parse_pptx(data: bytes, builder: CirBuilder) -> str:
    try:
        presentation = open_presentation(io.BytesIO(data))
    except (KeyError, OSError, ValueError, zipfile.BadZipFile) as exc:
        raise StructuredParseError("PPTX_PARSE_FAILED") from exc
    if len(presentation.slides) > builder.limits.max_slides:
        raise StructuredParseError("SLIDE_LIMIT")

    title = normalize_text(presentation.core_properties.title or "")
    slide_width = int(presentation.slide_width or 1)
    slide_height = int(presentation.slide_height or 1)
    slide_metadata: list[dict[str, Any]] = []
    for slide_index0, slide in enumerate(presentation.slides):
        flattened, group_count = _walk_shapes(slide.shapes)
        if len(flattened) > builder.limits.max_shapes_per_slide:
            raise StructuredParseError("SLIDE_SHAPE_LIMIT")
        title_shape = slide.shapes.title
        slide_title = (
            normalize_text(title_shape.text)
            if title_shape is not None and getattr(title_shape, "has_text_frame", False)
            else ""
        )
        if not slide_title:
            slide_title = f"Slide {slide_index0 + 1}"
        if not title and slide_index0 == 0:
            title = slide_title
        slide_location = SourceLocation(
            page_index0=slide_index0,
            native_object_id=f"pptx/slide/{slide_index0:04d}",
            bbox1000=(
                _shape_bbox(title_shape, slide_width, slide_height)
                if title_shape is not None
                else (0, 0, 1000, 1000)
            ),
        )
        slide_block = builder.add_block(
            block_type=BlockType.TITLE if slide_index0 == 0 else BlockType.HEADING,
            location=slide_location,
            raw_text=slide_title,
            markdown=f"# {slide_title}" if slide_index0 == 0 else f"## {slide_title}",
            quality_flags=(("slide_title_inferred",) if title_shape is None else ()),
        )

        reading_order, connector_metadata, reading_cycle = _reading_order(flattened)
        figure_count = 0
        table_count = 0
        emitted_order: list[str] = []
        for reading_index, node in enumerate(reading_order):
            if (
                title_shape is not None
                and getattr(node.shape, "shape_id", None) == title_shape.shape_id
            ):
                continue
            emitted = _add_shape(
                node,
                builder=builder,
                slide_index0=slide_index0,
                parent_id=slide_block.id,
                slide_width=slide_width,
                slide_height=slide_height,
                reading_index=reading_index,
            )
            if emitted is not None:
                emitted_order.append(f"pptx/slide/{slide_index0:04d}/shape/{node.z_path}")
                if emitted == BlockType.FIGURE:
                    figure_count += 1
                elif emitted == BlockType.TABLE:
                    table_count += 1

        note_text = ""
        if slide.has_notes_slide:
            note_text = normalize_text(slide.notes_slide.notes_text_frame.text)
        if note_text:
            builder.add_block(
                block_type=BlockType.FOOTNOTE,
                location=SourceLocation(
                    page_index0=slide_index0,
                    native_object_id=f"pptx/slide/{slide_index0:04d}/notes",
                ),
                raw_text=note_text,
                markdown=note_text,
                parent_id=slide_block.id,
                quality_flags=("speaker_notes",),
            )
        slide_metadata.append(
            {
                "pageIndex0": slide_index0,
                "shapeCount": len(flattened),
                "groupCount": group_count,
                "connectorCount": len(connector_metadata),
                "connectors": connector_metadata,
                "readingOrder": emitted_order,
                "readingOrderStrategy": "connector-topology-group-position-z",
                "readingOrderCycleDetected": reading_cycle,
                "figureCount": figure_count,
                "tableCount": table_count,
                "hasSpeakerNotes": bool(note_text),
            }
        )
    if not builder.blocks:
        raise StructuredParseError("PPTX_EMPTY_PRESENTATION")
    builder.metadata["slides"] = slide_metadata
    return title or _filename_title(builder.source_filename)


def _walk_shapes(shapes: Any) -> tuple[list[_ShapeNode], int]:
    nodes: list[_ShapeNode] = []
    group_count = 0

    def visit(
        collection: Any,
        *,
        prefix: tuple[int, ...] = (),
        group_path: str | None = None,
        container_top: int | None = None,
        container_left: int | None = None,
    ) -> None:
        nonlocal group_count
        for z_index, shape in enumerate(collection):
            path = (*prefix, z_index)
            encoded = ".".join(f"{value:04d}" for value in path)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                group_count += 1
                visit(
                    shape.shapes,
                    prefix=path,
                    group_path=group_path or encoded,
                    container_top=(
                        container_top
                        if container_top is not None
                        else int(getattr(shape, "top", 0))
                    ),
                    container_left=(
                        container_left
                        if container_left is not None
                        else int(getattr(shape, "left", 0))
                    ),
                )
                continue
            nodes.append(
                _ShapeNode(
                    z_path=encoded,
                    shape=shape,
                    group_path=group_path,
                    container_top=(
                        container_top
                        if container_top is not None
                        else int(getattr(shape, "top", 0))
                    ),
                    container_left=(
                        container_left
                        if container_left is not None
                        else int(getattr(shape, "left", 0))
                    ),
                )
            )

    visit(shapes)
    return nodes, group_count


def _reading_order(
    nodes: list[_ShapeNode],
) -> tuple[list[_ShapeNode], list[dict[str, Any]], bool]:
    by_shape_id = {
        int(node.shape.shape_id): node
        for node in nodes
        if getattr(node.shape, "shape_id", None) is not None
    }
    edges: dict[int, set[int]] = {shape_id: set() for shape_id in by_shape_id}
    incoming: dict[int, int] = {shape_id: 0 for shape_id in by_shape_id}
    connector_metadata: list[dict[str, Any]] = []
    connector_ids: set[int] = set()
    for node in nodes:
        start_id, end_id = _connector_endpoints(node.shape)
        if start_id is None and end_id is None:
            continue
        shape_id = int(getattr(node.shape, "shape_id", -1))
        connector_ids.add(shape_id)
        descriptor: dict[str, Any] = {"zPath": node.z_path}
        if start_id is not None:
            descriptor["startShapeId"] = start_id
        if end_id is not None:
            descriptor["endShapeId"] = end_id
        connector_metadata.append(descriptor)
        if (
            start_id is not None
            and end_id is not None
            and start_id in edges
            and end_id in edges
            and end_id not in edges[start_id]
            and start_id != end_id
        ):
            edges[start_id].add(end_id)
            incoming[end_id] += 1

    content_ids = set(by_shape_id) - connector_ids
    ready = sorted(
        (shape_id for shape_id in content_ids if incoming[shape_id] == 0),
        key=lambda shape_id: _shape_sort_key(by_shape_id[shape_id]),
    )
    ordered: list[_ShapeNode] = []
    while ready:
        shape_id = ready.pop(0)
        ordered.append(by_shape_id[shape_id])
        for target in sorted(
            edges[shape_id],
            key=lambda target_id: _shape_sort_key(by_shape_id[target_id]),
        ):
            incoming[target] -= 1
            if incoming[target] == 0 and target in content_ids:
                ready.append(target)
                ready.sort(key=lambda ready_id: _shape_sort_key(by_shape_id[ready_id]))

    remaining = content_ids - {
        int(node.shape.shape_id)
        for node in ordered
        if getattr(node.shape, "shape_id", None) is not None
    }
    cycle = bool(remaining)
    ordered.extend(
        by_shape_id[shape_id]
        for shape_id in sorted(
            remaining,
            key=lambda remaining_id: _shape_sort_key(by_shape_id[remaining_id]),
        )
    )
    unconnected = [node for node in nodes if getattr(node.shape, "shape_id", None) is None]
    ordered.extend(sorted(unconnected, key=_shape_sort_key))
    return ordered, connector_metadata, cycle


def _shape_sort_key(node: _ShapeNode) -> tuple[int, int, str, int, int, str]:
    return (
        node.container_top,
        node.container_left,
        node.group_path or node.z_path,
        int(getattr(node.shape, "top", 0)),
        int(getattr(node.shape, "left", 0)),
        node.z_path,
    )


def _connector_endpoints(shape: Any) -> tuple[int | None, int | None]:
    start_id: int | None = None
    end_id: int | None = None
    for element in shape._element.iter():
        local_name = str(element.tag).rsplit("}", 1)[-1]
        if local_name not in {"stCxn", "endCxn"}:
            continue
        raw_id = element.get("id")
        if raw_id is None or not str(raw_id).isdigit():
            continue
        if local_name == "stCxn":
            start_id = int(raw_id)
        else:
            end_id = int(raw_id)
    return start_id, end_id


def _add_shape(
    node: _ShapeNode,
    *,
    builder: CirBuilder,
    slide_index0: int,
    parent_id: str,
    slide_width: int,
    slide_height: int,
    reading_index: int,
) -> BlockType | None:
    shape = node.shape
    native_id = f"pptx/slide/{slide_index0:04d}/shape/{node.z_path}"
    location = SourceLocation(
        page_index0=slide_index0,
        native_object_id=native_id,
        bbox1000=_shape_bbox(shape, slide_width, slide_height),
    )
    if getattr(shape, "has_table", False):
        _add_pptx_table(
            shape.table,
            builder=builder,
            location=location,
            parent_id=parent_id,
        )
        return BlockType.TABLE

    if getattr(shape, "has_chart", False):
        _add_chart(
            shape,
            builder=builder,
            location=location,
            parent_id=parent_id,
            group_path=node.group_path,
            reading_index=reading_index,
        )
        return BlockType.FIGURE
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        _add_picture(
            shape,
            builder=builder,
            location=location,
            parent_id=parent_id,
            group_path=node.group_path,
            reading_index=reading_index,
        )
        return BlockType.FIGURE

    if not getattr(shape, "has_text_frame", False):
        return None
    emitted = False
    for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
        text = normalize_text(paragraph.text)
        if not text:
            continue
        paragraph_location = SourceLocation(
            page_index0=slide_index0,
            native_object_id=f"{native_id}/p/{paragraph_index:04d}",
            bbox1000=location.bbox1000,
        )
        block_type = _pptx_paragraph_type(shape, paragraph)
        flags = [
            f"shape_z_order:{node.z_path}",
            f"reading_order:{reading_index:04d}",
        ]
        if node.group_path:
            flags.append(f"shape_group:{node.group_path}")
        builder.add_block(
            block_type=block_type,
            location=paragraph_location,
            raw_text=text,
            markdown=_pptx_markdown(text, block_type),
            parent_id=parent_id,
            quality_flags=tuple(flags),
        )
        emitted = True
    return BlockType.PARAGRAPH if emitted else None


def _add_picture(
    shape: Any,
    *,
    builder: CirBuilder,
    location: SourceLocation,
    parent_id: str,
    group_path: str | None,
    reading_index: int,
) -> None:
    image = shape.image
    payload = bytes(image.blob)
    relationship_id = str(getattr(shape._element, "blip_rId", ""))
    image_part = shape.part.related_part(relationship_id) if relationship_id else None
    native_part = str(getattr(image_part, "partname", f"{location.native_object_id}/image")).lstrip(
        "/"
    )
    asset_id = builder.register_asset(
        payload=payload,
        native_part=native_part,
        media_type=str(getattr(image, "content_type", "application/octet-stream")),
        kind="image",
        filename=str(getattr(image, "filename", native_part.rsplit("/", 1)[-1])),
        width_px=int(image.size[0]),
        height_px=int(image.size[1]),
    )
    label, has_alt_text = _shape_label(shape)
    flags = [
        "embedded_asset_extracted",
        f"reading_order:{reading_index:04d}",
    ]
    if has_alt_text:
        flags.append("alt_text_preserved")
    if group_path:
        flags.append(f"shape_group:{group_path}")
    builder.add_block(
        block_type=BlockType.FIGURE,
        location=SourceLocation(
            page_index0=location.page_index0,
            native_object_id=location.native_object_id,
            bbox1000=location.bbox1000,
            image_asset_id=asset_id,
        ),
        raw_text=label,
        markdown=f"![{label}](akc-asset:{asset_id})",
        parent_id=parent_id,
        quality_flags=tuple(flags),
    )


def _add_chart(
    shape: Any,
    *,
    builder: CirBuilder,
    location: SourceLocation,
    parent_id: str,
    group_path: str | None,
    reading_index: int,
) -> None:
    chart = shape.chart
    chart_part = chart.part
    payload = bytes(chart_part.blob)
    chart_data = _chart_data(chart)
    native_part = str(chart_part.partname).lstrip("/")
    asset_id = builder.register_asset(
        payload=payload,
        native_part=native_part,
        media_type=str(chart_part.content_type),
        kind="chart",
        filename=native_part.rsplit("/", 1)[-1],
        metadata={
            "slidePlacement": location.native_object_id,
            "chartData": chart_data,
        },
    )
    label, has_alt_text = _shape_label(shape)
    chart_title = normalize_text(str(chart_data.get("title", "")))
    if chart_title:
        label = chart_title
    summary_lines = [label]
    for series in chart_data["series"]:
        series_name = str(series["name"])
        values = ", ".join(str(value) for value in series["values"])
        if values:
            summary_lines.append(f"{series_name}: {values}")
    flags = [
        "chart_structure_extracted",
        "embedded_asset_extracted",
        f"reading_order:{reading_index:04d}",
    ]
    if has_alt_text:
        flags.append("alt_text_preserved")
    if group_path:
        flags.append(f"shape_group:{group_path}")
    builder.add_block(
        block_type=BlockType.FIGURE,
        location=SourceLocation(
            page_index0=location.page_index0,
            native_object_id=location.native_object_id,
            bbox1000=location.bbox1000,
            image_asset_id=asset_id,
        ),
        raw_text="\n".join(summary_lines),
        markdown=f"**Chart: {label}**\n\n" + "\n".join(summary_lines[1:]),
        parent_id=parent_id,
        quality_flags=tuple(flags),
    )


def _chart_data(chart: Any) -> dict[str, Any]:
    title = ""
    if getattr(chart, "has_title", False):
        title = normalize_text(chart.chart_title.text_frame.text)
    series_data: list[dict[str, Any]] = []
    plot_types: list[str] = []
    for plot_index, plot in enumerate(chart.plots):
        plot_types.append(type(plot).__name__)
        categories = _chart_categories(plot)
        for series_index, series in enumerate(plot.series):
            name = normalize_text(str(getattr(series, "name", "")))
            if not name:
                name = f"Series {series_index + 1}"
            values = [_json_scalar(value) for value in tuple(getattr(series, "values", ()))]
            descriptor: dict[str, Any] = {
                "plotIndex0": plot_index,
                "seriesIndex0": series_index,
                "name": name,
                "values": values,
            }
            if categories:
                descriptor["categories"] = categories
            series_data.append(descriptor)
    return {
        "title": title,
        "plotTypes": plot_types,
        "series": series_data,
    }


def _chart_categories(plot: Any) -> list[str]:
    categories = getattr(plot, "categories", None)
    if categories is None:
        return []
    try:
        return [
            normalize_text(str(getattr(category, "label", category))) for category in categories
        ]
    except (AttributeError, TypeError, ValueError):
        return []


def _json_scalar(value: Any) -> str | int | float | bool | None:
    if isinstance(value, float) and not math.isfinite(value):
        return normalize_text(str(value))
    if value is None or isinstance(value, (str, int, float, bool)):
        return value
    return normalize_text(str(value))


def _pptx_paragraph_type(shape: Any, paragraph: Any) -> BlockType:
    if shape.is_placeholder:
        placeholder_type = shape.placeholder_format.type
        if placeholder_type in _TITLE_PLACEHOLDERS or placeholder_type == PP_PLACEHOLDER.SUBTITLE:
            return BlockType.HEADING
    paragraph_xml = paragraph._p.xml
    if paragraph.level > 0 or "<a:buChar" in paragraph_xml or "<a:buAutoNum" in paragraph_xml:
        return BlockType.LIST
    return BlockType.PARAGRAPH


def _pptx_markdown(text: str, block_type: BlockType) -> str:
    if block_type == BlockType.HEADING:
        return f"### {text}"
    if block_type == BlockType.LIST:
        return f"- {text}"
    return text


def _add_pptx_table(
    table: Any,
    *,
    builder: CirBuilder,
    location: SourceLocation,
    parent_id: str,
) -> None:
    row_count = len(table.rows)
    column_count = len(table.columns)
    cells: list[TableCellSpec] = []
    for row_index in range(row_count):
        for column_index in range(column_count):
            cell = table.cell(row_index, column_index)
            if cell.is_spanned:
                continue
            cells.append(
                TableCellSpec(
                    row_index0=row_index,
                    column_index0=column_index,
                    row_span=int(cell.span_height),
                    column_span=int(cell.span_width),
                    raw_text=normalize_text(cell.text),
                    location=SourceLocation(
                        page_index0=location.page_index0,
                        native_object_id=(
                            f"{location.native_object_id}/table/"
                            f"r/{row_index:04d}/c/{column_index:04d}"
                        ),
                        bbox1000=location.bbox1000,
                    ),
                )
            )
    if not cells:
        return
    header_row_count = 1 if all(cell.raw_text for cell in cells if cell.row_index0 == 0) else 0
    builder.add_table(
        location=SourceLocation(
            page_index0=location.page_index0,
            native_object_id=f"{location.native_object_id}/table",
            bbox1000=location.bbox1000,
        ),
        row_count=row_count,
        column_count=column_count,
        cells=tuple(cells),
        header_row_count=header_row_count,
        parent_id=parent_id,
        quality_flags=(("header_row_inferred",) if header_row_count else ()),
    )


def _shape_bbox(
    shape: Any,
    slide_width: int,
    slide_height: int,
) -> tuple[int, int, int, int] | None:
    if shape is None or slide_width <= 0 or slide_height <= 0:
        return None
    left = int(getattr(shape, "left", 0))
    top = int(getattr(shape, "top", 0))
    width = int(getattr(shape, "width", 0))
    height = int(getattr(shape, "height", 0))
    if width <= 0 or height <= 0:
        return None
    x1 = max(0, min(999, round(left / slide_width * 1000)))
    y1 = max(0, min(999, round(top / slide_height * 1000)))
    x2 = max(x1 + 1, min(1000, round((left + width) / slide_width * 1000)))
    y2 = max(y1 + 1, min(1000, round((top + height) / slide_height * 1000)))
    return x1, y1, x2, y2


def _shape_label(shape: Any) -> tuple[str, bool]:
    for element in shape._element.iter():
        if not str(element.tag).endswith("}cNvPr"):
            continue
        for attribute in ("descr", "title"):
            value = normalize_text(str(element.get(attribute, "")))
            if value:
                return value, True
        break
    return normalize_text(str(getattr(shape, "name", ""))) or "Slide figure", False


def _filename_title(filename: str) -> str:
    return normalize_text(filename.rsplit(".", 1)[0].replace("_", " ")) or "Presentation"
