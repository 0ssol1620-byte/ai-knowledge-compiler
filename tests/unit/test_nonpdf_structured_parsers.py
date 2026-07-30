from __future__ import annotations

import io
import zipfile
from collections.abc import Callable
from dataclasses import replace
from datetime import UTC, datetime
from pathlib import Path

import pytest
from akc_cir import BlockType, CanonicalDocument, canonical_json
from akc_native_parsers import (
    ParseContext,
    ParserLimits,
    StructuredParseError,
    parse_non_pdf_to_cir,
)
from docx import Document as WordDocument
from openpyxl import Workbook
from openpyxl.worksheet.table import Table as WorksheetTable
from pptx import Presentation
from pptx.util import Inches

FIXTURES = Path(__file__).parents[1] / "fixtures" / "nonpdf"
MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "html": "text/html",
    "srt": "application/x-subrip",
    "vtt": "text/vtt",
}


@pytest.fixture
def parse_context() -> ParseContext:
    return ParseContext(
        tenant_id="tenant_fixture",
        document_id="document_fixture",
        document_version_id="version_fixture",
        created_at=datetime(2026, 7, 29, 12, 0, tzinfo=UTC),
        source_url="https://ingest.example.test/source",
        retrieved_at=datetime(2026, 7, 29, 11, 59, tzinfo=UTC),
    )


@pytest.fixture
def docx_bytes() -> bytes:
    document = WordDocument()
    document.core_properties.title = "구조 보존 DOCX"
    document.add_paragraph("구조 보존 DOCX", style="Title")
    document.add_heading("실험 결과", level=1)
    document.add_paragraph("원문 문단입니다.")
    document.add_paragraph("첫 번째 항목", style="List Bullet")
    table = document.add_table(rows=2, cols=3)
    table.cell(0, 0).text = "구성"
    table.cell(0, 1).text = "점수"
    table.cell(0, 2).text = "비고"
    table.cell(1, 0).text = "개선"
    table.cell(1, 1).text = "0.94"
    table.cell(1, 1).merge(table.cell(1, 2))
    document.sections[0].header.paragraphs[0].text = "보안 등급: 내부"
    document.sections[0].footer.paragraphs[0].text = "문서 끝"
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "PPTX 구조"
    body = first.placeholders[1].text_frame
    body.text = "첫 번째 문단"
    bullet = body.add_paragraph()
    bullet.text = "핵심 항목"
    bullet.level = 1
    table = first.shapes.add_table(
        2,
        2,
        Inches(1),
        Inches(4),
        Inches(4),
        Inches(1.5),
    ).table
    table.cell(0, 0).text = "구성"
    table.cell(0, 1).text = "점수"
    table.cell(1, 0).text = "기준"
    table.cell(1, 1).text = "0.86"
    first.notes_slide.notes_text_frame.text = "발표자 전용 메모"
    second = deck.slides.add_slide(deck.slide_layouts[5])
    second.shapes.title.text = "두 번째 슬라이드"
    second.shapes.add_textbox(
        Inches(1),
        Inches(2),
        Inches(5),
        Inches(1),
    ).text_frame.text = "슬라이드 순서를 보존합니다."
    deck.core_properties.title = "구조 보존 PPTX"
    output = io.BytesIO()
    deck.save(output)
    return output.getvalue()


@pytest.fixture
def xlsx_bytes() -> bytes:
    workbook = Workbook()
    workbook.properties.title = "구조 보존 XLSX"
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Results"
    sheet.append(["구성", "점수", "비고", "병합", "합계"])
    sheet.append(["기준", 0.86, "검증 전", None, "=SUM(B2:B3)"])
    sheet.append(["개선", 0.94, "검증 후", None, None])
    sheet.merge_cells("C2:D2")
    sheet.row_dimensions[3].hidden = True
    sheet.add_table(WorksheetTable(displayName="EvidenceTable", ref="A1:B3"))
    hidden = workbook.create_sheet("Raw")
    hidden.sheet_state = "hidden"
    hidden.append(["raw", "value"])
    hidden.append(["alpha", 1])
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()
    return output.getvalue()


def _parse(
    *,
    filename: str,
    data: bytes,
    context: ParseContext,
    limits: ParserLimits | None = None,
) -> CanonicalDocument:
    extension = Path(filename).suffix.removeprefix(".")
    return parse_non_pdf_to_cir(
        filename=filename,
        declared_mime=MIME[extension],
        data=data,
        context=context,
        limits=limits,
    )


def test_docx_preserves_xml_order_hierarchy_table_and_story_parts(
    parse_context: ParseContext,
    docx_bytes: bytes,
) -> None:
    document = _parse(
        filename="evidence.docx",
        data=docx_bytes,
        context=parse_context,
    )
    types = [block.type for block in document.blocks]
    assert types[:4] == [
        BlockType.TITLE,
        BlockType.HEADING,
        BlockType.PARAGRAPH,
        BlockType.LIST,
    ]
    assert BlockType.TABLE in types
    assert BlockType.HEADER in types
    assert BlockType.FOOTER in types
    table_block = next(block for block in document.blocks if block.type == BlockType.TABLE)
    assert table_block.table is not None
    assert any(cell.column_span == 2 for cell in table_block.table.cells)
    assert all(
        (native_id := block.source_refs[0].native_object_id) is not None
        and native_id.startswith("docx/")
        for block in document.blocks
    )


def test_pptx_preserves_slide_shape_order_geometry_tables_and_notes(
    parse_context: ParseContext,
    pptx_bytes: bytes,
) -> None:
    document = _parse(
        filename="deck.pptx",
        data=pptx_bytes,
        context=parse_context,
    )
    assert document.metadata["slides"][0]["shapeCount"] >= 3
    assert {block.source_refs[0].page_index0 for block in document.blocks} == {0, 1}
    assert any(block.type == BlockType.TABLE for block in document.blocks)
    assert any(block.type == BlockType.LIST for block in document.blocks)
    assert any(block.type == BlockType.FOOTNOTE for block in document.blocks)
    positioned = [
        block
        for block in document.blocks
        if "/shape/" in (block.source_refs[0].native_object_id or "")
    ]
    assert positioned
    assert all(block.source_refs[0].bbox1000 is not None for block in positioned)


def test_xlsx_preserves_sheet_table_merge_formula_and_visibility(
    parse_context: ParseContext,
    xlsx_bytes: bytes,
) -> None:
    document = _parse(
        filename="workbook.xlsx",
        data=xlsx_bytes,
        context=parse_context,
    )
    assert [sheet["name"] for sheet in document.metadata["sheets"]] == [
        "Results",
        "Raw",
    ]
    assert document.metadata["sheets"][1]["state"] == "hidden"
    table_blocks = [block for block in document.blocks if block.type == BlockType.TABLE]
    assert len(table_blocks) == 2
    first_table = table_blocks[0].table
    assert first_table is not None
    assert any(cell.column_span == 2 for cell in first_table.cells)
    formula = next(cell for cell in first_table.cells if cell.raw_text.startswith("="))
    assert "formula_preserved_not_executed" in formula.quality_flags
    assert "formula_cached_value_missing" in formula.quality_flags
    assert document.metadata["sheets"][0]["tables"][0]["name"] == "EvidenceTable"


def test_html_preserves_dom_blocks_and_never_retains_active_fetches(
    parse_context: ParseContext,
) -> None:
    source = (FIXTURES / "sample.html").read_bytes()
    document = _parse(
        filename="article.html",
        data=source,
        context=parse_context,
    )
    types = {block.type for block in document.blocks}
    assert {
        BlockType.TITLE,
        BlockType.HEADING,
        BlockType.PARAGRAPH,
        BlockType.LIST,
        BlockType.TABLE,
        BlockType.FIGURE,
        BlockType.CAPTION,
    }.issubset(types)
    joined = "\n".join(block.raw_text or "" for block in document.blocks)
    assert "window.evil" not in joined
    assert "제거할 탐색" not in joined
    assert document.metadata["htmlSafety"]["strippedReferenceCount"] == 2
    assert "html_external_references_not_fetched" in document.metadata["warnings"]
    assert all(
        (native_id := block.source_refs[0].native_object_id) is not None
        and native_id.startswith("html/")
        for block in document.blocks
    )


def test_html_preserves_multiple_articles_and_unique_text_node_locations(
    parse_context: ParseContext,
) -> None:
    source = (
        b"<!doctype html><html><body>"
        b"<article>alpha <strong>beta</strong> omega<!--ignored--></article>"
        b"<article><p>second article</p></article>"
        b"</body></html>"
    )
    document = _parse(
        filename="articles.html",
        data=source,
        context=parse_context,
    )
    texts = [block.raw_text for block in document.blocks]
    assert texts == ["alpha", "beta", "omega", "second article"]
    assert len({block.id for block in document.blocks}) == len(document.blocks)
    locations = [block.source_refs[0].native_object_id for block in document.blocks]
    assert len(set(locations)) == len(locations)
    assert all("ignored" not in (text or "") for text in texts)


def test_subtitles_preserve_time_speaker_and_merge_repeated_cues(
    parse_context: ParseContext,
) -> None:
    srt = _parse(
        filename="captions.srt",
        data=(FIXTURES / "sample.srt").read_bytes(),
        context=parse_context,
    )
    assert srt.metadata["subtitles"]["sourceCueCount"] == 4
    assert srt.metadata["subtitles"]["canonicalCueCount"] == 3
    cue_blocks = [block for block in srt.blocks if block.type == BlockType.PARAGRAPH]
    assert len(cue_blocks) == 3
    assert cue_blocks[0].source_refs[0].time_start_ms == 1_000
    assert cue_blocks[0].source_refs[0].time_end_ms == 3_000
    assert "speaker_label_detected" in cue_blocks[0].quality_flags
    merged = next(block for block in cue_blocks if "repeated_cues_merged" in block.quality_flags)
    assert merged.source_refs[0].time_start_ms == 3_500
    assert merged.source_refs[0].time_end_ms == 6_500

    vtt = _parse(
        filename="captions.vtt",
        data=(FIXTURES / "sample.vtt").read_bytes(),
        context=parse_context,
    )
    assert vtt.metadata["subtitles"]["sourceCueCount"] == 2
    assert vtt.metadata["subtitles"]["skippedMetadataBlockCount"] == 2
    first_vtt_cue = next(block for block in vtt.blocks if block.type == BlockType.PARAGRAPH)
    assert first_vtt_cue.normalized_text == "첫 번째 관찰입니다."
    assert "speaker_label_detected" in first_vtt_cue.quality_flags


def test_deterministic_ids_order_and_wire_round_trip(
    parse_context: ParseContext,
    docx_bytes: bytes,
) -> None:
    first = _parse(
        filename="stable.docx",
        data=docx_bytes,
        context=parse_context,
    )
    second = _parse(
        filename="stable.docx",
        data=docx_bytes,
        context=parse_context,
    )
    assert first.model_dump(mode="json") == second.model_dump(mode="json")
    assert [block.order for block in first.blocks] == list(range(len(first.blocks)))
    wire = canonical_json(first)
    assert CanonicalDocument.model_validate_json(wire) == first


def test_all_real_minimal_samples_reach_canonical_cir(
    parse_context: ParseContext,
    docx_bytes: bytes,
    pptx_bytes: bytes,
    xlsx_bytes: bytes,
) -> None:
    sources = [
        ("sample.docx", docx_bytes),
        ("sample.pptx", pptx_bytes),
        ("sample.xlsx", xlsx_bytes),
        ("sample.html", (FIXTURES / "sample.html").read_bytes()),
        ("sample.srt", (FIXTURES / "sample.srt").read_bytes()),
        ("sample.vtt", (FIXTURES / "sample.vtt").read_bytes()),
    ]
    for filename, payload in sources:
        document = _parse(
            filename=filename,
            data=payload,
            context=parse_context,
        )
        assert document.blocks
        assert all(block.source_refs for block in document.blocks)
        assert CanonicalDocument.model_validate_json(canonical_json(document)) == document


def test_fail_closed_mime_macro_external_relation_and_embedding(
    parse_context: ParseContext,
    docx_bytes: bytes,
) -> None:
    with pytest.raises(StructuredParseError, match="MIME_MISMATCH"):
        parse_non_pdf_to_cir(
            filename="sample.docx",
            declared_mime="application/pdf",
            data=docx_bytes,
            context=parse_context,
        )
    with pytest.raises(StructuredParseError, match="UNSUPPORTED_NON_PDF_TYPE"):
        parse_non_pdf_to_cir(
            filename="sample.docm",
            declared_mime="application/octet-stream",
            data=docx_bytes,
            context=parse_context,
        )

    external = _rewrite_zip(
        docx_bytes,
        transform={
            "word/_rels/document.xml.rels": lambda value: value.replace(
                b"</Relationships>",
                (
                    b'<Relationship Id="unsafe" '
                    b'Type="http://schemas.openxmlformats.org/officeDocument/'
                    b'2006/relationships/hyperlink" '
                    b'Target="https://example.invalid/" TargetMode="External"/>'
                    b"</Relationships>"
                ),
            )
        },
    )
    with pytest.raises(StructuredParseError, match="OFFICE_EXTERNAL_RELATION"):
        _parse(
            filename="external.docx",
            data=external,
            context=parse_context,
        )

    embedded = _rewrite_zip(
        docx_bytes,
        additions={"word/embeddings/object.bin": b"untrusted object"},
    )
    with pytest.raises(StructuredParseError, match="OFFICE_EMBEDDED_OBJECT"):
        _parse(
            filename="embedded.docx",
            data=embedded,
            context=parse_context,
        )

    macro_enabled = _rewrite_zip(
        docx_bytes,
        transform={
            "[Content_Types].xml": lambda value: value.replace(
                b"</Types>",
                (
                    b'<Override PartName="/word/macro.bin" '
                    b'ContentType="application/vnd.ms-office.'
                    b'vbaProjectMacroEnabled.main+xml"/></Types>'
                ),
            )
        },
    )
    with pytest.raises(StructuredParseError, match="OFFICE_ACTIVE_CONTENT"):
        _parse(
            filename="macro-disguised.docx",
            data=macro_enabled,
            context=parse_context,
        )


def test_format_specific_limits_fail_before_unbounded_work(
    parse_context: ParseContext,
    pptx_bytes: bytes,
    xlsx_bytes: bytes,
) -> None:
    with pytest.raises(StructuredParseError, match="SLIDE_LIMIT"):
        _parse(
            filename="deck.pptx",
            data=pptx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_slides=1),
        )
    with pytest.raises(StructuredParseError, match="SHEET_ROW_LIMIT"):
        _parse(
            filename="book.xlsx",
            data=xlsx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_rows_per_sheet=2),
        )
    with pytest.raises(StructuredParseError, match="SUBTITLE_CUE_LIMIT"):
        _parse(
            filename="captions.srt",
            data=(FIXTURES / "sample.srt").read_bytes(),
            context=parse_context,
            limits=replace(ParserLimits(), max_subtitle_cues=2),
        )
    with pytest.raises(StructuredParseError, match="HTML_NODE_LIMIT"):
        _parse(
            filename="article.html",
            data=(FIXTURES / "sample.html").read_bytes(),
            context=parse_context,
            limits=replace(ParserLimits(), max_html_nodes=5),
        )


def test_archive_size_ratio_magic_column_and_table_area_limits_fail_closed(
    parse_context: ParseContext,
    docx_bytes: bytes,
    xlsx_bytes: bytes,
) -> None:
    with pytest.raises(StructuredParseError, match="FILE_TOO_LARGE"):
        _parse(
            filename="sample.docx",
            data=docx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_input_bytes=len(docx_bytes) - 1),
        )
    with pytest.raises(StructuredParseError, match="ARCHIVE_ENTRY_LIMIT"):
        _parse(
            filename="sample.docx",
            data=docx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_archive_entries=1),
        )
    with pytest.raises(StructuredParseError, match="ARCHIVE_SIZE_LIMIT"):
        _parse(
            filename="sample.docx",
            data=docx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_archive_uncompressed_bytes=1_000),
        )
    with pytest.raises(StructuredParseError, match="ARCHIVE_RATIO_LIMIT"):
        _parse(
            filename="sample.docx",
            data=docx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_compression_ratio=1.01),
        )
    with pytest.raises(StructuredParseError, match="MAGIC_MISMATCH"):
        _parse(
            filename="fake.docx",
            data=b"<!doctype html><html><body>not office</body></html>",
            context=parse_context,
        )
    with pytest.raises(StructuredParseError, match="SHEET_COLUMN_LIMIT"):
        _parse(
            filename="wide.xlsx",
            data=xlsx_bytes,
            context=parse_context,
            limits=replace(ParserLimits(), max_columns_per_sheet=4),
        )
    with pytest.raises(StructuredParseError, match="TABLE_CELL_LIMIT"):
        _parse(
            filename="span.html",
            data=(
                b"<!doctype html><html><body><table><tr>"
                b'<td rowspan="100000" colspan="6">bounded</td>'
                b"</tr></table></body></html>"
            ),
            context=parse_context,
        )


def test_malformed_ooxml_is_reported_as_stable_parse_failure(
    parse_context: ParseContext,
    docx_bytes: bytes,
) -> None:
    malformed = _rewrite_zip(
        docx_bytes,
        transform={"word/document.xml": lambda _value: b"<not-valid-xml"},
    )
    with pytest.raises(StructuredParseError) as failure:
        _parse(
            filename="malformed.docx",
            data=malformed,
            context=parse_context,
        )
    assert failure.value.code == "DOCX_PARSE_FAILED"


def _rewrite_zip(
    source: bytes,
    *,
    transform: dict[str, Callable[[bytes], bytes]] | None = None,
    additions: dict[str, bytes] | None = None,
) -> bytes:
    transforms = transform or {}
    output = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(source)) as archive,
        zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as rewritten,
    ):
        for entry in archive.infolist():
            payload = archive.read(entry)
            callback = transforms.get(entry.filename)
            if callback is not None:
                payload = callback(payload)
            rewritten.writestr(entry, payload)
        for name, payload in (additions or {}).items():
            rewritten.writestr(name, payload)
    return output.getvalue()
