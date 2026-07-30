from __future__ import annotations

import io
import zipfile
from collections.abc import Callable
from datetime import UTC, datetime

import pytest
from akc_cir import BlockType
from akc_native_parsers import (
    ParseContext,
    StructuredParseError,
    parse_non_pdf_to_cir,
)
from docx import Document as WordDocument
from lxml import etree
from openpyxl import Workbook
from openpyxl.chart import BarChart, Reference
from openpyxl.drawing.image import Image as SpreadsheetImage
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches

_WORD_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"
_MIME = {
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "srt": "application/x-subrip",
}


def _context() -> ParseContext:
    return ParseContext(
        tenant_id="tenant_fidelity",
        document_id="document_fidelity",
        document_version_id="version_fidelity",
        created_at=datetime(2026, 7, 30, 1, 0, tzinfo=UTC),
    )


def _parse(filename: str, payload: bytes):
    extension = filename.rsplit(".", 1)[-1]
    return parse_non_pdf_to_cir(
        filename=filename,
        declared_mime=_MIME[extension],
        data=payload,
        context=_context(),
    )


def _png_bytes() -> bytes:
    output = io.BytesIO()
    Image.new("RGB", (24, 12), color=(33, 96, 144)).save(output, format="PNG")
    return output.getvalue()


def _docx_with_assets_comments_and_revisions() -> bytes:
    document = WordDocument()
    document.add_paragraph("Native fidelity", style="Title")
    image_paragraph = document.add_paragraph()
    run = image_paragraph.add_run()
    run.add_picture(io.BytesIO(_png_bytes()), width=Inches(1))
    doc_properties = run._r.xpath(".//wp:docPr")[0]
    doc_properties.set("descr", "Architecture overview")
    document.add_paragraph("Figure 1. Architecture", style="Caption")
    commented = document.add_paragraph("Review this statement")
    document.add_comment(
        runs=commented.runs,
        text="Verified reviewer note",
        author="Reviewer",
        initials="RV",
    )
    document.add_paragraph("revision target")
    document.sections[0].header.paragraphs[0].text = "Confidential"
    document.sections[0].footer.paragraphs[0].text = "Page footer"
    output = io.BytesIO()
    document.save(output)

    def add_revisions(value: bytes) -> bytes:
        root = etree.fromstring(value)
        namespace = {"w": _WORD_NS}
        target = root.xpath("//w:t[text()='revision target']", namespaces=namespace)[0]
        run_element = target.getparent()
        paragraph = run_element.getparent()
        target.text = "base "
        insertion = etree.Element(f"{{{_WORD_NS}}}ins")
        insertion.set(f"{{{_WORD_NS}}}author", "Editor")
        insertion_run = etree.SubElement(insertion, f"{{{_WORD_NS}}}r")
        insertion_text = etree.SubElement(insertion_run, f"{{{_WORD_NS}}}t")
        insertion_text.text = "visible insertion"
        deletion = etree.Element(f"{{{_WORD_NS}}}del")
        deletion.set(f"{{{_WORD_NS}}}author", "Editor")
        deletion_run = etree.SubElement(deletion, f"{{{_WORD_NS}}}r")
        deletion_text = etree.SubElement(deletion_run, f"{{{_WORD_NS}}}delText")
        deletion_text.text = "deleted text"
        index = paragraph.index(run_element)
        paragraph.insert(index + 1, insertion)
        paragraph.insert(index + 2, deletion)
        return etree.tostring(
            root,
            xml_declaration=True,
            encoding="UTF-8",
            standalone=True,
        )

    return _rewrite_zip(
        output.getvalue(),
        transform={"word/document.xml": add_revisions},
    )


def test_docx_extracts_embedded_asset_caption_comment_and_revision_views() -> None:
    document = _parse(
        "fidelity.docx",
        _docx_with_assets_comments_and_revisions(),
    )
    assets = document.metadata["assets"]
    assert len(assets) == 1
    assert assets[0]["kind"] == "image"
    assert assets[0]["mediaType"] == "image/png"
    assert assets[0]["widthPx"] == 24
    assert assets[0]["heightPx"] == 12
    assert assets[0]["sha256"].startswith("sha256:")
    figure = next(block for block in document.blocks if block.type == BlockType.FIGURE)
    assert figure.source_refs[0].image_asset_id == assets[0]["id"]
    assert "alt_text_preserved" in figure.quality_flags
    caption = next(block for block in document.blocks if block.type == BlockType.CAPTION)
    assert caption.parent_id == figure.id
    comment = next(block for block in document.blocks if "docx_comment" in block.quality_flags)
    assert comment.raw_text == "Verified reviewer note"
    assert comment.parent_id is not None
    visible = "\n".join(block.raw_text or "" for block in document.blocks)
    assert "base visible insertion" in visible
    assert "deleted text" not in visible
    revisions = document.metadata["docx"]["trackedChanges"]
    assert [revision["kind"] for revision in revisions] == ["insertion", "deletion"]
    assert revisions[0]["visible"] is True
    assert revisions[1]["text"] == "deleted text"
    assert document.metadata["docx"]["trackedChangeView"] == (
        "insertions-visible-deletions-metadata-only"
    )
    assert any(block.type == BlockType.HEADER for block in document.blocks)
    assert any(block.type == BlockType.FOOTER for block in document.blocks)


def _pptx_with_assets_chart_group_and_connector() -> bytes:
    deck = Presentation()
    slide = deck.slides.add_slide(deck.slide_layouts[5])
    slide.shapes.title.text = "Visual evidence"
    picture = slide.shapes.add_picture(
        io.BytesIO(_png_bytes()),
        Inches(0.5),
        Inches(1.5),
        Inches(1.5),
        Inches(0.75),
    )
    picture._pic.nvPicPr.cNvPr.set("descr", "System architecture")

    chart_data = ChartData()
    chart_data.categories = ["Baseline", "Candidate"]
    chart_data.add_series("Accuracy", (0.82, 0.94))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(3),
        Inches(1.5),
        Inches(4),
        Inches(2.5),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Accuracy"

    group = slide.shapes.add_group_shape()
    grouped_first = group.shapes.add_textbox(
        Inches(1),
        Inches(4.5),
        Inches(2),
        Inches(0.5),
    )
    grouped_first.text_frame.text = "Grouped source"
    grouped_second = group.shapes.add_textbox(
        Inches(4),
        Inches(4.5),
        Inches(2),
        Inches(0.5),
    )
    grouped_second.text_frame.text = "Grouped target"
    connector = group.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(3),
        Inches(4.7),
        Inches(4),
        Inches(4.7),
    )
    connector.begin_connect(grouped_first, 3)
    connector.end_connect(grouped_second, 1)
    slide.notes_slide.notes_text_frame.text = "Speaker-only context"
    output = io.BytesIO()
    deck.save(output)
    return output.getvalue()


def test_pptx_extracts_image_chart_notes_and_group_connector_reading_order() -> None:
    document = _parse(
        "visuals.pptx",
        _pptx_with_assets_chart_group_and_connector(),
    )
    assets = document.metadata["assets"]
    assert {asset["kind"] for asset in assets} == {"image", "chart"}
    picture = next(
        block
        for block in document.blocks
        if block.type == BlockType.FIGURE and "alt_text_preserved" in block.quality_flags
    )
    assert picture.source_refs[0].image_asset_id is not None
    chart = next(
        block for block in document.blocks if "chart_structure_extracted" in block.quality_flags
    )
    assert "Accuracy" in (chart.raw_text or "")
    chart_asset = next(asset for asset in assets if asset["kind"] == "chart")
    chart_data = chart_asset["metadata"]["chartData"]
    assert chart_data["series"][0]["name"] == "Accuracy"
    assert chart_data["series"][0]["values"] == [0.82, 0.94]
    slide = document.metadata["slides"][0]
    assert slide["groupCount"] == 1
    assert slide["connectorCount"] == 1
    assert slide["readingOrderStrategy"] == "connector-topology-group-position-z"
    assert any("shape_group:" in flag for block in document.blocks for flag in block.quality_flags)
    assert any(
        block.type == BlockType.FOOTNOTE and "speaker_notes" in block.quality_flags
        for block in document.blocks
    )


def test_pptx_chart_workbook_is_relationship_scoped_and_recursively_validated() -> None:
    source = _pptx_with_assets_chart_group_and_connector()
    with zipfile.ZipFile(io.BytesIO(source)) as archive:
        embedding_name = next(
            name
            for name in archive.namelist()
            if name.casefold().startswith("ppt/embeddings/") and name.casefold().endswith(".xlsx")
        )
        embedded_workbook = archive.read(embedding_name)

    unrelated = _rewrite_zip(
        source,
        additions={"ppt/embeddings/unrelated.xlsx": embedded_workbook},
    )
    with pytest.raises(StructuredParseError, match="OFFICE_EMBEDDED_OBJECT"):
        _parse("unrelated.pptx", unrelated)

    def add_external_relationship(value: bytes) -> bytes:
        return value.replace(
            b"</Relationships>",
            (
                b'<Relationship Id="unsafe" '
                b'Type="http://schemas.openxmlformats.org/officeDocument/'
                b'2006/relationships/externalLink" '
                b'Target="https://example.invalid/" TargetMode="External"/>'
                b"</Relationships>"
            ),
        )

    unsafe_workbook = _rewrite_zip(
        embedded_workbook,
        transform={"xl/_rels/workbook.xml.rels": add_external_relationship},
    )
    unsafe_chart = _rewrite_zip(
        source,
        transform={embedding_name: lambda _value: unsafe_workbook},
    )
    with pytest.raises(StructuredParseError, match="OFFICE_EXTERNAL_RELATION"):
        _parse("unsafe-chart.pptx", unsafe_chart)


def _xlsx_with_assets_chart_formula_and_hidden_state() -> bytes:
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Evidence"
    sheet.append(["Category", "Value", "Formula"])
    sheet.append(["Baseline", 1, "=SUM(B2:B3)"])
    sheet.append(["Candidate", 2, None])
    sheet.row_dimensions[3].hidden = True
    sheet.column_dimensions["B"].hidden = True
    image = SpreadsheetImage(io.BytesIO(_png_bytes()))
    image.anchor = "E2"
    sheet.add_image(image)
    chart = BarChart()
    chart.title = "Comparison"
    chart.add_data(Reference(sheet, min_col=2, min_row=1, max_row=3), titles_from_data=True)
    chart.set_categories(Reference(sheet, min_col=1, min_row=2, max_row=3))
    sheet.add_chart(chart, "E8")
    output = io.BytesIO()
    workbook.save(output)
    workbook.close()

    def add_formula_cache(value: bytes) -> bytes:
        return value.replace(
            b"<f>SUM(B2:B3)</f><v></v>",
            b"<f>SUM(B2:B3)</f><v>3</v>",
        )

    return _rewrite_zip(
        output.getvalue(),
        transform={"xl/worksheets/sheet1.xml": add_formula_cache},
    )


def test_xlsx_extracts_assets_chart_formula_cache_and_hidden_state() -> None:
    document = _parse(
        "evidence.xlsx",
        _xlsx_with_assets_chart_formula_and_hidden_state(),
    )
    sheet = document.metadata["sheets"][0]
    assert sheet["hiddenRows"] == [3]
    assert sheet["hiddenColumns"] == ["B"]
    assert sheet["formulas"] == [
        {
            "cell": "C2",
            "formula": "=SUM(B2:B3)",
            "cachedValue": "3",
            "cachedValuePresent": True,
        }
    ]
    assert sheet["imageCount"] == 1
    assert sheet["chartCount"] == 1
    assets = document.metadata["assets"]
    assert {asset["kind"] for asset in assets} == {"image", "chart"}
    assert all(asset["sha256"].startswith("sha256:") for asset in assets)
    figures = [block for block in document.blocks if block.type == BlockType.FIGURE]
    assert len(figures) == 2
    assert all(block.source_refs[0].image_asset_id for block in figures)
    chart_asset = next(asset for asset in assets if asset["kind"] == "chart")
    assert chart_asset["metadata"]["chartData"]["series"][0]["valueReference"].endswith("$B$2:$B$3")


def test_subtitle_segments_preserve_cues_with_deterministic_topic_boundaries() -> None:
    payload = (
        b"1\n00:00:00,000 --> 00:00:20,000\nOpening context.\n\n"
        b"2\n00:00:20,000 --> 00:00:40,000\nSupporting detail.\n\n"
        b"3\n00:00:40,000 --> 00:01:05,000\nFirst topic closes.\n\n"
        b"4\n00:01:10,000 --> 00:01:20,000\nNew topic starts.\n\n"
        b"5\n00:01:20,000 --> 00:01:45,000\nNew topic develops.\n"
    )
    document = _parse("topics.srt", payload)
    segmentation = document.metadata["subtitles"]["segmentation"]
    assert segmentation["segmentCount"] == 2
    assert [segment["durationMs"] for segment in segmentation["segments"]] == [
        65_000,
        35_000,
    ]
    assert all(segment["within30To90Seconds"] for segment in segmentation["segments"])
    segment_blocks = [
        block for block in document.blocks if "subtitle_segment" in block.quality_flags
    ]
    cue_blocks = [block for block in document.blocks if block.type == BlockType.PARAGRAPH]
    assert len(segment_blocks) == 2
    assert len(cue_blocks) == 5
    assert {block.parent_id for block in cue_blocks} == {block.id for block in segment_blocks}
    assert segment_blocks[0].source_refs[0].time_start_ms == 0
    assert segment_blocks[-1].source_refs[0].time_end_ms == 105_000


def _rewrite_zip(
    source: bytes,
    *,
    transform: dict[str, Callable[[bytes], bytes]] | None = None,
    additions: dict[str, bytes] | None = None,
) -> bytes:
    transforms = transform or {}
    output = io.BytesIO()
    with (
        zipfile.ZipFile(io.BytesIO(source)) as archive,
        zipfile.ZipFile(output, "w", compression=zipfile.ZIP_DEFLATED) as rewritten,
    ):
        for entry in archive.infolist():
            payload = archive.read(entry)
            callback = transforms.get(entry.filename)
            if callback is not None:
                payload = callback(payload)
            rewritten.writestr(entry, payload)
        for name, payload in (additions or {}).items():
            rewritten.writestr(name, payload)
    return output.getvalue()
